#!/usr/bin/env python3
"""
Fetch Sumo Logic virtual training schedule and produce:
  1. training_schedule.html  — filterable web page (future sessions shown by default)
  2. training_schedule.pptx  — compact slide deck suitable for presentations

Timezones displayed: PT, NZ (Auckland), AET (Sydney), SGT (Singapore).

Requirements:
    pip install playwright beautifulsoup4 python-pptx
    playwright install chromium
"""

import json
import re
import sys
from datetime import datetime, timezone
from pathlib import Path
from zoneinfo import ZoneInfo

try:
    from playwright.sync_api import sync_playwright
except ImportError:
    print("ERROR: playwright not installed.")
    print("Run:  pip install playwright && playwright install chromium")
    sys.exit(1)

from bs4 import BeautifulSoup

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------
TIMEZONES = [
    ("PT",  "America/Los_Angeles", "Los Angeles"),
    ("NZ",  "Pacific/Auckland",    "Auckland"),
    ("AET", "Australia/Sydney",    "Sydney"),
    ("SGT", "Asia/Singapore",      "Singapore"),
]

OUTPUT_HTML = Path(__file__).parent / "training_schedule.html"
OUTPUT_PPTX = Path(__file__).parent / "training_schedule.pptx"
TARGET_URL  = "https://www.sumologic.com/learn/training"

# ---------------------------------------------------------------------------
# Fetch
# ---------------------------------------------------------------------------

def fetch_all_pages(url: str) -> list[str]:
    """
    Load the page, detect FacetWP pagination, click through every page,
    and return a list of HTML strings (one per page).
    """
    print(f"Fetching {url} ...")
    pages_html = []

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(
            user_agent=(
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/124.0.0.0 Safari/537.36"
            )
        )
        page.goto(url, wait_until="load", timeout=90_000)

        # Wait for FacetWP to complete its initial AJAX population.
        # networkidle gives us 500 ms of no network activity, which covers the
        # FacetWP AJAX call that fires after DOMContentLoaded.
        try:
            page.wait_for_load_state("networkidle", timeout=20_000)
        except Exception:
            pass  # non-fatal; we'll rely on the item-count check below

        # Ensure at least some cards are in the DOM before we read the pager
        try:
            page.wait_for_selector(".blog-item", state="attached", timeout=15_000)
        except Exception:
            print("  WARNING: No .blog-item cards found after waiting.")

        # Scroll to trigger any lazy-load / infinite-scroll then let it settle
        page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
        try:
            page.wait_for_load_state("networkidle", timeout=10_000)
        except Exception:
            page.wait_for_timeout(2000)

        def _max_page_from_pager() -> int:
            """
            Parse the facetwp-pager and return the highest page number visible.
            FacetWP pagers often only render a window of links (e.g. 1 2 3 … 8),
            so we also check for a 'last' page link that carries a higher data-page.
            """
            if not page.query_selector(".facetwp-pager"):
                return 1
            pager_html = page.inner_html(".facetwp-pager")
            soup_p = BeautifulSoup(pager_html, "html.parser")
            nums = set()
            for a in soup_p.find_all("a", class_="facetwp-page"):
                dp = a.get("data-page", "")
                if dp.isdigit():
                    nums.add(int(dp))
            # Also check span/div elements that FacetWP uses for the current/last page
            for el in soup_p.find_all(True):
                dp = el.get("data-page", "")
                if dp.isdigit():
                    nums.add(int(dp))
            return max(nums) if nums else 1

        max_page = _max_page_from_pager()
        item_count_p1 = page.evaluate("document.querySelectorAll('.blog-item').length")
        print(f"  Detected {max_page} page(s) via facetwp-pager. "
              f"Page 1 initial load: {item_count_p1} .blog-item cards.")

        # Workaround: sumologic.com has a site-side pagination bug where the first load
        # of page 1 renders an incomplete set of cards. Navigating to page 2 then back
        # forces FacetWP to re-render page 1 with the correct full list.
        if max_page > 1:
            print("  Applying page-1 refresh workaround (site pagination bug)...")
            page.click(".facetwp-pager a.facetwp-page[data-page='2']")
            try:
                page.wait_for_load_state("networkidle", timeout=15_000)
            except Exception:
                page.wait_for_timeout(2000)
            page.click(".facetwp-pager a.facetwp-page[data-page='1']")
            try:
                page.wait_for_load_state("networkidle", timeout=15_000)
            except Exception:
                page.wait_for_timeout(2000)
            item_count_p1 = page.evaluate("document.querySelectorAll('.blog-item').length")
            print(f"  After workaround: page 1 has {item_count_p1} .blog-item cards.")

        # Collect page 1
        pages_html.append(page.content())

        # Click through pages 2..N
        pg = 2
        while pg <= max_page:
            print(f"  Loading page {pg}/{max_page} ...")
            selector = f".facetwp-pager a.facetwp-page[data-page='{pg}']"
            # Snapshot item count before click so we can detect content replacement
            prev_count = page.evaluate("document.querySelectorAll('.blog-item').length")
            page.click(selector)
            # Wait for FacetWP to swap in new content (count must change, then stabilise)
            try:
                page.wait_for_function(
                    f"() => document.querySelectorAll('.blog-item').length !== {prev_count}",
                    timeout=20_000,
                )
            except Exception:
                # Fall back: wait for the loading spinner to clear if present
                try:
                    page.wait_for_selector(".facetwp-loading", state="detached", timeout=10_000)
                except Exception:
                    pass
            page.wait_for_timeout(1000)
            pages_html.append(page.content())
            # Re-read the pager — later pages may reveal a higher total (e.g. "…" resolved)
            new_max = _max_page_from_pager()
            if new_max > max_page:
                print(f"  Pager updated: max page is now {new_max}.")
                max_page = new_max
            pg += 1

        browser.close()

    return pages_html


# ---------------------------------------------------------------------------
# Parse schedule
# ---------------------------------------------------------------------------

def parse_schedule(soup: BeautifulSoup) -> list[dict]:
    """
    Parse instructor-led virtual class cards.
    Each card is a <div class="blog-item"> inside an <a> that links to registration.
    Times are ISO datetimes in <time data-start="..." data-end="...">.
    """
    sessions = []

    # Find the section containing the schedule
    schedule_section = None
    for h in soup.find_all(["h2", "h3"]):
        if re.search(r"instructor.led", h.get_text(), re.I):
            schedule_section = h.find_parent("section")
            break

    if not schedule_section:
        print("  WARNING: Could not find 'Instructor-Led' section container.")
        return sessions

    # Each card is: <a class="blog-item-link" href="..."><div class="blog-item">...</div></a>
    # Search the whole document for blog-item-link anchors (they may sit outside the section tag)
    search_root = schedule_section or soup
    for a_tag in soup.find_all("a", class_="blog-item-link", href=True):
        card = a_tag.find("div", class_="blog-item")
        if not card:
            continue
        title_el = card.find("h3")
        time_el  = card.find("time", attrs={"data-start": True})
        if not title_el or not time_el:
            continue

        title      = title_el.get_text(strip=True)
        data_start = time_el.get("data-start", "")
        data_end   = time_el.get("data-end", "")
        reg_url    = a_tag["href"]

        sessions.append({
            "title":      title,
            "data_start": data_start,
            "data_end":   data_end,
            "reg_url":    reg_url,
        })

    print(f"  Found {len(sessions)} sessions.")
    return sessions


# ---------------------------------------------------------------------------
# Parse certifications descriptions
# ---------------------------------------------------------------------------

def parse_descriptions(soup: BeautifulSoup) -> dict[str, str]:
    """
    Return {course_title: description_text} from the Certifications overview section.
    """
    descs = {}
    for block in soup.find_all("div", class_="resource-block__body-wrap"):
        heading = block.find(re.compile(r"^h[2-6]$"))
        para    = block.find("p")
        if heading and para:
            descs[heading.get_text(strip=True)] = para.get_text(strip=True)
    print(f"  Found {len(descs)} certification descriptions.")
    return descs


# ---------------------------------------------------------------------------
# Timezone conversion
# ---------------------------------------------------------------------------

def convert_times(sessions: list[dict]) -> list[dict]:
    """
    Add a 'tz_columns' dict and an ISO sort key to each session.
    """
    for s in sessions:
        dt = None
        if s["data_start"]:
            try:
                # Python's fromisoformat handles offsets like -07:00
                dt = datetime.fromisoformat(s["data_start"])
            except ValueError:
                pass

        s["dt"] = dt
        s["sort_key"] = dt.isoformat() if dt else "9999"
        s["tz_columns"] = {}

        for label, tz_name, _ in TIMEZONES:
            if dt:
                local_dt  = dt.astimezone(ZoneInfo(tz_name))
                s["tz_columns"][label] = local_dt.strftime("%a %d %b %Y  %I:%M %p")
            else:
                s["tz_columns"][label] = s["data_start"]

        # End time in PT for display
        if dt and s["data_end"]:
            try:
                dt_end = datetime.fromisoformat(s["data_end"])
                pt_end = dt_end.astimezone(ZoneInfo("America/Los_Angeles"))
                s["pt_end"] = pt_end.strftime("%I:%M %p")
            except ValueError:
                s["pt_end"] = ""
        else:
            s["pt_end"] = ""

    sessions.sort(key=lambda x: x["sort_key"])
    return sessions


# ---------------------------------------------------------------------------
# Best-match description lookup
# ---------------------------------------------------------------------------

KNOWN_MAPPINGS = {
    "welcome webinar":           "Fundamentals",
    "fundamental":               "Fundamentals",
    "fundamentals":              "Fundamentals",
    "search mastery":            "Search Mastery",
    "best practice":             "Search Mastery",
    "query efficiency":          "Search Mastery",
    "admin":                     "Administration",
    "administration":            "Administration",
    "metrics":                   "Metrics",
    "logs for security":         "Logs for Security",
    "cloud siem":                "Cloud SIEM",
    "automation":                "Cloud SIEM",
    "mobot":                     "Search Mastery",
}


def find_description(title: str, descs: dict[str, str]) -> str:
    title_lower = title.lower()
    # Direct keyword match
    for keyword, cert_title in KNOWN_MAPPINGS.items():
        if keyword in title_lower and cert_title in descs:
            return descs[cert_title]
    # Fallback: overlap scoring
    title_words = set(re.sub(r"[^a-z0-9 ]", "", title_lower).split())
    best, best_score = "", 0
    for cert_title, desc in descs.items():
        cert_words = set(re.sub(r"[^a-z0-9 ]", "", cert_title.lower()).split())
        score = len(title_words & cert_words)
        if score > best_score:
            best, best_score = desc, score
    return best


# ---------------------------------------------------------------------------
# Title normalisation
# ---------------------------------------------------------------------------

_AMPM_RE = re.compile(r'\s+[AP]M$', re.IGNORECASE)


def strip_ampm(title: str) -> str:
    """Remove trailing AM / PM slot indicator from a course title."""
    return _AMPM_RE.sub('', title).strip()


# ---------------------------------------------------------------------------
# HTML output
# ---------------------------------------------------------------------------

HTML_TEMPLATE = """\
<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width, initial-scale=1.0"/>
<title>Sumo Logic Virtual Training Schedule</title>
<style>
  :root {{
    --bg:#0f1117; --card:#1a1d27; --border:#2e3348;
    --accent:#00b4d8; --accent2:#ff6b35; --text:#e2e8f0;
    --muted:#94a3b8; --green:#22c55e; --purple:#a855f7;
  }}
  *{{ box-sizing:border-box; margin:0; padding:0; }}
  body{{ background:var(--bg); color:var(--text);
    font-family:'Segoe UI',system-ui,sans-serif; min-height:100vh; padding:24px; }}

  header{{ display:flex; align-items:center; gap:16px;
    margin-bottom:28px; padding-bottom:16px; border-bottom:1px solid var(--border); }}
  h1{{ font-size:1.55rem; color:var(--accent); }}
  .subtitle{{ color:var(--muted); font-size:0.88rem; margin-top:4px; }}
  .source-link{{ color:var(--accent); font-size:0.82rem; text-decoration:none; }}
  .source-link:hover{{ text-decoration:underline; }}

  .controls{{
    display:flex; flex-wrap:wrap; gap:10px; margin-bottom:20px; align-items:center;
  }}
  .controls input,.controls select{{
    background:var(--card); border:1px solid var(--border); color:var(--text);
    padding:8px 13px; border-radius:8px; font-size:0.88rem; outline:none;
    transition:border-color .2s;
  }}
  .controls input:focus,.controls select:focus{{ border-color:var(--accent); }}
  .controls input{{ width:250px; }}
  .count{{ color:var(--muted); font-size:0.83rem; margin-left:auto; }}

  .table-wrap{{ overflow-x:auto; border-radius:12px; border:1px solid var(--border); }}
  table{{ width:100%; border-collapse:collapse; font-size:0.87rem; }}
  thead tr{{ background:var(--card); border-bottom:2px solid var(--border); }}
  th{{
    padding:11px 13px; text-align:left; color:var(--accent);
    font-weight:600; white-space:nowrap; cursor:pointer; user-select:none;
  }}
  th:hover{{ color:var(--text); }}
  th .si{{ opacity:.4; margin-left:4px; font-size:.72rem; }}
  th.sorted .si{{ opacity:1; }}
  tbody tr{{ border-bottom:1px solid var(--border); transition:background .15s; }}
  tbody tr:hover{{ background:rgba(0,180,216,.06); }}
  tbody tr.hidden{{ display:none; }}
  td{{ padding:10px 13px; vertical-align:top; }}
  td.title-cell{{ min-width:200px; max-width:280px; }}

  .course-btn{{
    background:none; border:none; color:var(--accent); font-weight:600;
    font-size:0.87rem; cursor:pointer; text-align:left; padding:0;
    display:flex; align-items:flex-start; gap:6px; width:100%;
  }}
  .course-btn:hover{{ color:var(--text); }}
  .ei{{ font-size:.68rem; color:var(--muted); transition:transform .2s; margin-top:3px; flex-shrink:0; }}

  .desc-row{{ display:none; }}
  .desc-row.open{{ display:table-row; }}
  .desc-cell{{
    padding:6px 13px 14px 32px;
    background:rgba(0,180,216,.04);
    color:var(--muted); font-size:.84rem; line-height:1.65;
  }}
  .desc-cell strong{{ color:var(--text); display:block; margin-bottom:4px; font-size:.82rem; }}

  .reg-btn{{
    display:inline-block; padding:5px 13px; background:var(--accent2);
    color:#fff; border-radius:6px; font-size:.81rem; text-decoration:none;
    font-weight:600; white-space:nowrap;
  }}
  .reg-btn:hover{{ opacity:.85; }}
  .tba{{ color:var(--muted); font-size:.8rem; }}

  .tz-val{{ font-size:.82rem; white-space:nowrap; color:var(--text); }}
  .tz-end{{ font-size:.75rem; color:var(--muted); }}

  .toggle-label{{
    display:flex; align-items:center; gap:7px;
    color:var(--muted); font-size:.87rem; cursor:pointer; user-select:none;
  }}
  .toggle-label input{{ accent-color:var(--accent); width:15px; height:15px; cursor:pointer; }}
  .export-btn{{
    background:var(--card); border:1px solid var(--border); color:var(--accent);
    padding:8px 14px; border-radius:8px; font-size:.87rem; cursor:pointer;
    font-weight:600; transition:border-color .2s, color .2s;
  }}
  .export-btn:hover{{ border-color:var(--accent); color:var(--text); }}
  .export-btn.copied{{ border-color:var(--green); color:var(--green); }}
  tr.past td{{ opacity:.45; }}
  .no-results{{
    text-align:center; padding:48px; color:var(--muted); display:none;
  }}
  .summary-bar{{
    display:flex; flex-wrap:wrap; gap:8px; margin-bottom:16px;
    padding:10px 14px; background:var(--card); border:1px solid var(--border);
    border-radius:10px; align-items:center;
  }}
  .sum-label{{ color:var(--muted); font-size:.8rem; margin-right:2px; flex-shrink:0; }}
  .sum-badges{{ display:flex; flex-wrap:wrap; gap:6px; }}
  .sum-badge{{
    background:rgba(0,180,216,.1); border:1px solid rgba(0,180,216,.2);
    color:var(--text); font-size:.77rem; padding:3px 10px; border-radius:20px;
    cursor:pointer; transition:background .15s; white-space:nowrap;
  }}
  .sum-badge:hover{{ background:rgba(0,180,216,.3); }}
  .sum-badge .cnt{{ color:var(--accent); font-weight:700; margin-left:5px; }}
  footer{{
    margin-top:28px; text-align:center; color:var(--muted); font-size:.79rem;
  }}
  footer a{{ color:var(--accent); text-decoration:none; }}
</style>
</head>
<body>
<header>
  <div>
    <h1>Sumo Logic — Virtual Training Schedule</h1>
    <div class="subtitle">
      Instructor-Led Virtual Classes &nbsp;·&nbsp; Generated {fetch_date}
      &nbsp;·&nbsp;
      <a class="source-link" href="https://www.sumologic.com/learn/training" target="_blank">
        sumologic.com/learn/training ↗
      </a>
    </div>
  </div>
</header>

<div class="controls">
  <input type="text" id="searchBox" placeholder="Filter course… (comma or OR for multiple)" oninput="applyFilters()"/>
  <select id="monthFilter" onchange="applyFilters()">
    <option value="">All months</option>
  </select>
  <label class="toggle-label">
    <input type="checkbox" id="showPast" onchange="applyFilters()"/>
    Show past sessions
  </label>
  <span class="count" id="countLabel"></span>
  <button class="export-btn" id="exportBtn" onclick="exportToClipboard()">
    Copy for email
  </button>
</div>

<div class="summary-bar" id="summaryBar">
  <span class="sum-label">Courses:</span>
  <span class="sum-badges" id="summaryContent"></span>
</div>

<div class="table-wrap">
<table id="scheduleTable">
  <thead>
    <tr>
      <th onclick="sortTable(0)">Course <span class="si">⇅</span></th>
      <th onclick="sortTable(1)">🇺🇸 PT — Los Angeles <span class="si">⇅</span></th>
      <th onclick="sortTable(2)">🇳🇿 NZ — Auckland <span class="si">⇅</span></th>
      <th onclick="sortTable(3)">🇦🇺 AET — Sydney <span class="si">⇅</span></th>
      <th onclick="sortTable(4)">🇸🇬 SGT — Singapore <span class="si">⇅</span></th>
      <th>Register</th>
    </tr>
  </thead>
  <tbody id="tableBody">
{table_rows}
  </tbody>
</table>
</div>
<div class="no-results" id="noResults">No sessions match your filter.</div>

<footer>
  Data sourced from
  <a href="https://www.sumologic.com/learn/training" target="_blank">sumologic.com/learn/training</a>
  &nbsp;·&nbsp; Generated by <code>fetch_training_schedule.py</code>
</footer>

<script>
const SESSIONS = {sessions_json};

// Populate month filter from PT column values
const months = [];
const monthSet = new Set();
SESSIONS.forEach(s => {{
  if (s.pt) {{
    // e.g. "Tue 07 Apr 2026  05:00 AM"
    const parts = s.pt.split(/\\s+/);
    if (parts.length >= 4) {{
      const m = parts[2] + ' ' + parts[3];   // "Apr 2026"
      if (!monthSet.has(m)) {{ monthSet.add(m); months.push(m); }}
    }}
  }}
}});
months.forEach(m => {{
  const opt = document.createElement('option');
  opt.value = m; opt.textContent = m;
  document.getElementById('monthFilter').appendChild(opt);
}});

const now = Date.now();

function applyFilters() {{
  const raw      = document.getElementById('searchBox').value.toLowerCase().trim();
  // Support OR: split on comma or the word "or" (case-insensitive)
  const terms    = raw ? raw.split(/,|\\bor\\b/i).map(t => t.trim()).filter(Boolean) : [];
  const mon      = document.getElementById('monthFilter').value;
  const showPast = document.getElementById('showPast').checked;
  let visible = 0, total = 0;
  const courseCounts = {{}};
  document.querySelectorAll('tbody tr.data-row').forEach(row => {{
    const title    = row.dataset.title.toLowerCase();
    const dispTitle = row.dataset.title;   // already stripped of AM/PM
    const pt       = row.dataset.pt || '';
    const isoStr   = row.dataset.iso || '';
    const isFuture = !isoStr || (new Date(isoStr).getTime() >= now);
    row.classList.toggle('past', !isFuture);
    const matchesSearch = !terms.length || terms.some(t => title.includes(t));
    const show = (showPast || isFuture) && matchesSearch && (!mon || pt.includes(mon));
    row.classList.toggle('hidden', !show);
    if (!show) {{
      const dr = document.getElementById('desc-' + row.dataset.idx);
      if (dr) dr.classList.remove('open');
    }}
    if (show) {{
      visible++;
      courseCounts[dispTitle] = (courseCounts[dispTitle] || 0) + 1;
    }}
    total++;
  }});
  document.getElementById('countLabel').textContent =
    `Showing ${{visible}} of ${{total}} sessions`;
  document.getElementById('noResults').style.display = visible === 0 ? 'block' : 'none';
  // Rebuild course summary badges (sorted by count desc, then alpha)
  const summaryEl = document.getElementById('summaryContent');
  if (summaryEl) {{
    const entries = Object.entries(courseCounts)
      .sort((a, b) => b[1] - a[1] || a[0].localeCompare(b[0]));
    summaryEl.innerHTML = entries.map(([name, cnt]) =>
      `<span class="sum-badge" onclick="filterByCourse(${{JSON.stringify(name)}})">${{name}}<span class="cnt">${{cnt}}</span></span>`
    ).join('');
  }}
}}

function filterByCourse(name) {{
  document.getElementById('searchBox').value = name;
  applyFilters();
}}

function toggleDesc(idx) {{
  const dr   = document.getElementById('desc-' + idx);
  const icon = document.getElementById('ei-' + idx);
  if (!dr) return;
  const open = dr.classList.toggle('open');
  if (icon) icon.style.transform = open ? 'rotate(90deg)' : '';
}}

let sortCol = -1, sortAsc = true;
function sortTable(col) {{
  const tbody = document.getElementById('tableBody');
  const rows  = [...tbody.querySelectorAll('tr.data-row')];
  sortAsc = sortCol === col ? !sortAsc : true;
  sortCol = col;
  document.querySelectorAll('th').forEach((th, i) => {{
    th.classList.toggle('sorted', i === col);
    const si = th.querySelector('.si');
    if (si) si.textContent = i === col ? (sortAsc ? '↑' : '↓') : '⇅';
  }});
  rows.sort((a, b) => {{
    // Column 0: use data-title (clean name without the ▶ expand arrow)
    // Date columns: use data-iso on the row for chronological accuracy
    let av, bv;
    if (col === 0) {{
      av = a.dataset.title || '';
      bv = b.dataset.title || '';
    }} else if (col >= 1 && col <= 4) {{
      av = a.dataset.iso || a.querySelectorAll('td')[col]?.textContent.trim() || '';
      bv = b.dataset.iso || b.querySelectorAll('td')[col]?.textContent.trim() || '';
    }} else {{
      av = a.querySelectorAll('td')[col]?.textContent.trim() || '';
      bv = b.querySelectorAll('td')[col]?.textContent.trim() || '';
    }}
    return sortAsc ? av.localeCompare(bv) : bv.localeCompare(av);
  }});
  rows.forEach(row => {{
    const dr = document.getElementById('desc-' + row.dataset.idx);
    tbody.appendChild(row);
    if (dr) tbody.appendChild(dr);
  }});
}}

applyFilters();

// ---------------------------------------------------------------------------
// Export visible rows as HTML table (pastes as formatted table in email)
// ---------------------------------------------------------------------------
function exportToClipboard() {{
  const visibleRows = [...document.querySelectorAll('tbody tr.data-row')]
    .filter(r => !r.classList.contains('hidden'));

  if (!visibleRows.length) {{
    alert('No sessions are currently visible to export.');
    return;
  }}

  // ---- Build plain-text version (fallback) ----
  const COL_WIDTHS = [32, 26, 26, 26, 26];
  const HEADERS    = ['Course', 'PT (Los Angeles)', 'NZ (Auckland)', 'AET (Sydney)', 'SGT (Singapore)'];
  const pad = (s, n) => s.substring(0, n).padEnd(n);
  const divider = HEADERS.map((_, i) => '-'.repeat(COL_WIDTHS[i])).join('  ');

  let plain = 'Sumo Logic — Virtual Training Schedule\\n';
  plain += 'Register: https://www.sumologic.com/learn/training\\n\\n';
  plain += HEADERS.map((h, i) => pad(h, COL_WIDTHS[i])).join('  ') + '\\n';
  plain += divider + '\\n';

  visibleRows.forEach(row => {{
    const cells = row.querySelectorAll('td');
    const vals  = [0,1,2,3,4].map(i => {{
      const t = cells[i]?.innerText.replace(/\\n+/g, ' ').trim() || '';
      // strip the expand arrow prefix if present
      return t.replace(/^[▶▼]\\s*/, '');
    }});
    plain += vals.map((v, i) => pad(v, COL_WIDTHS[i])).join('  ') + '\\n';
    // Add registration URL on a sub-line under the course name
    const regLink = row.querySelector('td:last-child a');
    if (regLink) plain += pad('  ' + regLink.href, COL_WIDTHS[0]) + '\\n';
  }});

  // ---- Build HTML version (renders as table in Gmail / Outlook) ----
  const fontStack = "'Segoe UI', Arial, sans-serif";
  const colourHdr = "#00b4d8";
  const colourBg1 = "#f8fafc";
  const colourBg2 = "#ffffff";
  const colourBdr = "#cbd5e1";

  let html = `
<table style="border-collapse:collapse;font-family:${{fontStack}};font-size:13px;width:100%;">
  <thead>
    <tr style="background:#0f1117;">
      ${{HEADERS.map(h =>
        `<th style="padding:8px 10px;text-align:left;color:${{colourHdr}};
                    border:1px solid ${{colourBdr}};white-space:nowrap;">${{h}}</th>`
      ).join('')}}
      <th style="padding:8px 10px;text-align:left;color:${{colourHdr}};
                 border:1px solid ${{colourBdr}};">Register</th>
    </tr>
  </thead>
  <tbody>`;

  visibleRows.forEach((row, idx) => {{
    const cells  = row.querySelectorAll('td');
    const bgCol  = idx % 2 === 0 ? colourBg1 : colourBg2;
    const tdBase = `style="padding:7px 10px;border:1px solid ${{colourBdr}};
                           vertical-align:top;background:${{bgCol}};"`;
    const regLink = row.querySelector('td:last-child a');
    const regCell = regLink
      ? `<a href="${{regLink.href}}" style="color:#ff6b35;font-weight:600;">Register →</a>`
      : '<span style="color:#94a3b8;">TBA</span>';

    const colVals = [0,1,2,3,4].map(i => {{
      const t = cells[i]?.innerText.replace(/^[▶▼]\\s*/, '').replace(/\\n+/g, '<br>').trim() || '';
      return t;
    }});

    html += `
    <tr>
      ${{colVals.map(v => `<td ${{tdBase}}>${{v}}</td>`).join('')}}
      <td ${{tdBase}}>${{regCell}}</td>
    </tr>`;
  }});

  html += `
  </tbody>
</table>
<p style="font-size:11px;color:#64748b;margin-top:6px;">
  Source: <a href="https://www.sumologic.com/learn/training">sumologic.com/learn/training</a>
  &nbsp;·&nbsp; Generated {fetch_date}
</p>`;

  // ---- Write both to clipboard ----
  const btn = document.getElementById('exportBtn');
  if (navigator.clipboard && window.ClipboardItem) {{
    const item = new ClipboardItem({{
      'text/html':  new Blob([html],  {{ type: 'text/html' }}),
      'text/plain': new Blob([plain], {{ type: 'text/plain' }}),
    }});
    navigator.clipboard.write([item]).then(() => {{
      btn.textContent = '✓ Copied!';
      btn.classList.add('copied');
      setTimeout(() => {{ btn.textContent = 'Copy for email'; btn.classList.remove('copied'); }}, 2200);
    }}).catch(() => fallbackCopy(plain, btn));
  }} else {{
    fallbackCopy(plain, btn);
  }}
}}

function fallbackCopy(text, btn) {{
  // textarea trick for older browsers
  const ta = document.createElement('textarea');
  ta.value = text;
  ta.style.position = 'fixed'; ta.style.opacity = '0';
  document.body.appendChild(ta);
  ta.select();
  document.execCommand('copy');
  document.body.removeChild(ta);
  btn.textContent = '✓ Copied (plain text)';
  btn.classList.add('copied');
  setTimeout(() => {{ btn.textContent = 'Copy for email'; btn.classList.remove('copied'); }}, 2200);
}}
</script>
</body>
</html>
"""


def build_rows(sessions: list[dict], descs: dict[str, str]) -> tuple[str, list]:
    html_parts = []
    js_data    = []

    for i, s in enumerate(sessions):
        title         = s["title"]
        display_title = strip_ampm(title)
        reg_url       = s.get("reg_url", "#")
        cols          = s.get("tz_columns", {})
        pt_val        = cols.get("PT", "")
        nz_val        = cols.get("NZ", "")
        aet_val       = cols.get("AET", "")
        sgt_val       = cols.get("SGT", "")
        iso_start     = s.get("data_start", "")
        pt_end        = s.get("pt_end", "")

        desc = find_description(display_title, descs)

        # Title cell — clickable if description exists
        if desc:
            title_html = (
                f'<button class="course-btn" onclick="toggleDesc({i})">'
                f'<span class="ei" id="ei-{i}">▶</span>{display_title}</button>'
            )
        else:
            title_html = f'<span class="course-btn" style="cursor:default">{display_title}</span>'

        # PT end time annotation
        pt_display = pt_val
        if pt_end:
            pt_display = f'{pt_val}<br/><span class="tz-end">ends {pt_end}</span>'

        # Register
        if reg_url and reg_url != "#":
            reg_html = f'<a class="reg-btn" href="{reg_url}" target="_blank">Register →</a>'
        else:
            reg_html = '<span class="tba">TBA</span>'

        html_parts.append(
            f'    <tr class="data-row" data-idx="{i}" '
            f'data-title="{display_title}" data-pt="{pt_val}" data-iso="{iso_start}">\n'
            f'      <td class="title-cell">{title_html}</td>\n'
            f'      <td class="tz-val">{pt_display}</td>\n'
            f'      <td class="tz-val">{nz_val}</td>\n'
            f'      <td class="tz-val">{aet_val}</td>\n'
            f'      <td class="tz-val">{sgt_val}</td>\n'
            f'      <td>{reg_html}</td>\n'
            f'    </tr>\n'
        )

        # Description expand row
        if desc:
            desc_body = f"<strong>{display_title} — course overview</strong>{desc}"
        else:
            desc_body = "<em>No description available.</em>"

        html_parts.append(
            f'    <tr class="desc-row" id="desc-{i}">\n'
            f'      <td class="desc-cell" colspan="6">{desc_body}</td>\n'
            f'    </tr>\n'
        )

        js_data.append({"idx": i, "title": display_title, "pt": pt_val})

    return "".join(html_parts), js_data


def save_html(sessions: list[dict], descs: dict[str, str]) -> None:
    rows_html, js_data = build_rows(sessions, descs)
    html = HTML_TEMPLATE.format(
        fetch_date    = datetime.now().strftime("%d %b %Y %H:%M"),
        table_rows    = rows_html,
        sessions_json = json.dumps(js_data, ensure_ascii=False),
    )
    OUTPUT_HTML.write_text(html, encoding="utf-8")
    print(f"  Saved: {OUTPUT_HTML}")


# ---------------------------------------------------------------------------
# PowerPoint output
# ---------------------------------------------------------------------------

# Sumo Logic brand colours
SUMO_DARK   = "0F1117"   # near-black background
SUMO_NAVY   = "1A1D27"   # card/header fill
SUMO_TEAL   = "00B4D8"   # accent
SUMO_ORANGE = "FF6B35"   # secondary accent
SUMO_TEXT   = "E2E8F0"   # body text
SUMO_MUTED  = "94A3B8"   # secondary text
SUMO_BORDER = "2E3348"   # table grid colour

ROWS_PER_SLIDE = 10      # data rows per slide (excluding header)


def _rgb(hex6: str):
    """Convert 6-char hex string to pptx RGBColor."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))


def _set_cell(cell, text: str, *, bold=False, font_size=9,
              fg=SUMO_TEXT, bg=None, align="left", wrap=True):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    cell.text = text
    tf = cell.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER}.get(align, PP_ALIGN.LEFT)
    run = para.runs[0] if para.runs else para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(fg)

    if bg:
        from pptx.oxml.ns import qn
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgbClr  = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", bg)


def _add_slide(prs, title_text: str, page_num: int, total_pages: int):
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    blank_layout = prs.slide_layouts[6]   # completely blank
    slide = prs.slides.add_slide(blank_layout)

    # Background rectangle
    bg = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(SUMO_DARK)
    bg.line.fill.background()

    # Header bar
    hdr = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.6),
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = _rgb(SUMO_NAVY)
    hdr.line.fill.background()

    # Title text
    txb = slide.shapes.add_textbox(Inches(0.18), Inches(0.08), Inches(9), Inches(0.46))
    tf  = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_text
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = _rgb(SUMO_TEAL)

    # Page number (top-right)
    pgb = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.4), Inches(0.08), Inches(1.2), Inches(0.46),
    )
    ptf = pgb.text_frame
    pp = ptf.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    pr = pp.add_run()
    pr.text = f"p {page_num}/{total_pages}"
    pr.font.size = Pt(9)
    pr.font.color.rgb = _rgb(SUMO_MUTED)

    return slide


def save_pptx(sessions: list[dict]) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from lxml import etree

    # Filter to future sessions (same as HTML default view)
    now_utc = datetime.now(timezone.utc)
    future  = [s for s in sessions if s.get("dt") and s["dt"] >= now_utc]
    # Sort by datetime
    future.sort(key=lambda x: x["sort_key"])

    if not future:
        future = sessions   # fallback: show all if nothing is future

    # Split into pages
    chunks = [future[i:i + ROWS_PER_SLIDE] for i in range(0, len(future), ROWS_PER_SLIDE)]
    total_pages = len(chunks)

    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Column layout  [label, width_inches, wrap]
    COLS = [
        ("Course",              2.5,  True),
        ("PT  (Los Angeles)",   2.1,  False),
        ("NZ  (Auckland)",      2.1,  False),
        ("AET  (Sydney)",       2.1,  False),
        ("SGT  (Singapore)",    2.1,  False),
        ("Register URL",        2.0,  True),
    ]
    table_left   = Inches(0.18)
    table_top    = Inches(0.7)
    table_width  = sum(Inches(c[1]) for c in COLS)
    header_h     = Inches(0.34)
    row_h        = Inches(0.46)

    title_base = "Sumo Logic — Virtual Training Schedule"

    for page_idx, chunk in enumerate(chunks):
        page_num   = page_idx + 1

        # Date range label from the sessions on this slide
        dts = [s["dt"] for s in chunk if s.get("dt")]
        if dts:
            first = min(dts).astimezone(ZoneInfo("America/Los_Angeles"))
            last  = max(dts).astimezone(ZoneInfo("America/Los_Angeles"))
            if first.strftime("%b %Y") == last.strftime("%b %Y"):
                date_range = first.strftime("%d %b") + " – " + last.strftime("%d %b %Y")
            else:
                date_range = first.strftime("%d %b %Y") + " – " + last.strftime("%d %b %Y")
        else:
            date_range = datetime.now().strftime("%d %b %Y")

        title_text = (
            f"{title_base}   ·   {date_range}"
            if total_pages == 1
            else f"{title_base}   ·   {date_range}   ({page_num}/{total_pages})"
        )
        slide      = _add_slide(prs, title_text, page_num, total_pages)
        n_rows     = len(chunk) + 1   # +1 for header

        table_height = header_h + row_h * len(chunk)
        tbl_shape = slide.shapes.add_table(
            n_rows, len(COLS),
            table_left, table_top,
            table_width, table_height,
        )
        tbl = tbl_shape.table

        # Set column widths
        for ci, (_, w, _) in enumerate(COLS):
            tbl.columns[ci].width = Inches(w)

        # Set all row heights
        tbl.rows[0].height = header_h
        for ri in range(1, n_rows):
            tbl.rows[ri].height = row_h

        # Header row
        headers = [c[0] for c in COLS]
        for ci, hdr_text in enumerate(headers):
            _set_cell(
                tbl.cell(0, ci), hdr_text,
                bold=True, font_size=9, fg=SUMO_TEAL, bg=SUMO_NAVY,
                align="center",
            )

        # Data rows
        for ri, s in enumerate(chunk, start=1):
            cols    = s.get("tz_columns", {})
            pt_val  = cols.get("PT", "")
            pt_end  = s.get("pt_end", "")
            pt_disp = f"{pt_val}  →  {pt_end}" if pt_end else pt_val

            row_data = [
                s["title"],
                pt_disp,
                cols.get("NZ",  ""),
                cols.get("AET", ""),
                cols.get("SGT", ""),
                s.get("reg_url", ""),
            ]
            row_bg = SUMO_NAVY if ri % 2 == 0 else SUMO_DARK

            for ci, (cell_text, (_, _, wrap)) in enumerate(zip(row_data, COLS)):
                fsize = 7 if ci == 5 else 8   # smaller font for URL column
                _set_cell(
                    tbl.cell(ri, ci), cell_text,
                    font_size=fsize, fg=SUMO_TEXT, bg=row_bg, wrap=wrap,
                )

        # Apply grid borders to all cells
        for ri in range(n_rows):
            for ci in range(len(COLS)):
                tc   = tbl.cell(ri, ci)._tc
                tcPr = tc.get_or_add_tcPr()
                for side in ("lnL", "lnR", "lnT", "lnB"):
                    ln = etree.SubElement(tcPr, qn(f"a:{side}"))
                    ln.set("w", "6350")   # 0.5pt in EMUs
                    ln.set("cap", "flat")
                    solidFill = etree.SubElement(ln, qn("a:solidFill"))
                    srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
                    srgbClr.set("val", SUMO_BORDER)

        # Footer note
        note_left = table_left
        note_top  = table_top + table_height + Inches(0.12)
        noteb = slide.shapes.add_textbox(note_left, note_top, table_width, Inches(0.3))
        ntf  = noteb.text_frame
        np_  = ntf.paragraphs[0]
        nr   = np_.add_run()
        nr.text = "Register links are live URLs — click to open in browser.  All times shown as start time; PT column also shows end time."
        nr.font.size = Pt(7)
        nr.font.color.rgb = _rgb(SUMO_MUTED)
        nr.font.italic = True

    prs.save(str(OUTPUT_PPTX))
    print(f"  Saved: {OUTPUT_PPTX}")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    pages_html = fetch_all_pages(TARGET_URL)

    # Save debug copy of page 1
    debug = Path(__file__).parent / "debug_raw.html"
    debug.write_text(pages_html[0], encoding="utf-8")

    # Parse descriptions from page 1 (they're always there)
    soup1 = BeautifulSoup(pages_html[0], "html.parser")
    print("Parsing certification descriptions ...")
    descs = parse_descriptions(soup1)

    # Collect sessions from all pages
    print("Parsing schedule (all pages) ...")
    sessions = []
    for i, html in enumerate(pages_html, 1):
        soup = BeautifulSoup(html, "html.parser")
        page_sessions = parse_schedule(soup)
        print(f"  Page {i}: {len(page_sessions)} sessions")
        sessions.extend(page_sessions)
    print(f"  Total: {len(sessions)} sessions across {len(pages_html)} page(s).")

    if not sessions:
        print("\nWARNING: No sessions found. See debug_raw.html for the raw page source.")

    sessions = convert_times(sessions)

    print("Generating HTML ...")
    save_html(sessions, descs)

    print("Generating PowerPoint ...")
    save_pptx(sessions)

    if sessions:
        now_utc = datetime.now(timezone.utc)
        future_count = sum(1 for s in sessions if s.get("dt") and s["dt"] >= now_utc)
        print(f"\nDone! {len(sessions)} total sessions ({future_count} upcoming).")
        print(f"  HTML  → {OUTPUT_HTML.name}  (future sessions shown by default)")
        print(f"  PPTX  → {OUTPUT_PPTX.name}  ({future_count} upcoming sessions, date-sorted)")


if __name__ == "__main__":
    main()
